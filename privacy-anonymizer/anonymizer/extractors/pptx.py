"""Extractor de texto para presentaciones PowerPoint (.pptx)."""

from io import BytesIO
from typing import List

from pptx import Presentation

from anonymizer.extractors.base import BaseExtractor, ExtractedText


class PptxExtractor(BaseExtractor):
    """Extrae texto de archivos PowerPoint, incluyendo slides, notas y masters."""

    def _extract_from_shape(self, shape, slide_idx: int, shape_idx: int, extracted: List[ExtractedText]) -> None:
        """Extrae texto recursivamente de un shape, incluyendo grupos."""
        if hasattr(shape, "shapes"):
            # Grupo de shapes
            for sub_idx, sub_shape in enumerate(shape.shapes):
                self._extract_from_shape(sub_shape, slide_idx, f"{shape_idx}_group_{sub_idx}", extracted)
            return

        if not hasattr(shape, "text_frame"):
            return

        text_frame = shape.text_frame
        for para_idx, paragraph in enumerate(text_frame.paragraphs):
            for run_idx, run in enumerate(paragraph.runs):
                if run.text:
                    extracted.append(
                        ExtractedText(
                            text=run.text,
                            location=("slide", slide_idx, shape_idx, para_idx, run_idx),
                        )
                    )

    def extract(self, file_stream: BytesIO) -> List[ExtractedText]:
        prs = Presentation(file_stream)
        extracted: List[ExtractedText] = []

        # Slides
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                self._extract_from_shape(shape, slide_idx, shape_idx, extracted)

            # Notas del slide
            if slide.has_notes_slide and slide.notes_slide:
                notes_text_frame = slide.notes_slide.notes_text_frame
                for para_idx, paragraph in enumerate(notes_text_frame.paragraphs):
                    for run_idx, run in enumerate(paragraph.runs):
                        if run.text:
                            extracted.append(
                                ExtractedText(
                                    text=run.text,
                                    location=("notes", slide_idx, para_idx, run_idx),
                                )
                            )

        return extracted

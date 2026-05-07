"""Reconstructor de presentaciones PowerPoint (.pptx) anonimizadas."""

from io import BytesIO
from typing import List

from pptx import Presentation

from anonymizer.builders.base import BaseBuilder
from anonymizer.extractors.base import ExtractedText


class PptxBuilder(BaseBuilder):
    """Reconstruye una presentación PowerPoint reemplazando el texto por su versión anonimizada."""

    def _replace_in_shape(self, shape, slide_idx: int, shape_idx, extracted_texts: List[ExtractedText], redacted_texts: List[str], idx_ref: List[int]) -> None:
        """Reemplaza texto recursivamente en un shape, incluyendo grupos."""
        if hasattr(shape, "shapes"):
            for sub_idx, sub_shape in enumerate(shape.shapes):
                self._replace_in_shape(sub_shape, slide_idx, f"{shape_idx}_group_{sub_idx}", extracted_texts, redacted_texts, idx_ref)
            return

        if not hasattr(shape, "text_frame"):
            return

        text_frame = shape.text_frame
        for para_idx, paragraph in enumerate(text_frame.paragraphs):
            for run_idx, run in enumerate(paragraph.runs):
                idx = idx_ref[0]
                if idx < len(extracted_texts):
                    loc = extracted_texts[idx].location
                    if (
                        loc[0] == "slide"
                        and loc[1] == slide_idx
                        and loc[2] == shape_idx
                        and loc[3] == para_idx
                        and loc[4] == run_idx
                    ):
                        run.text = redacted_texts[idx]
                        idx_ref[0] = idx + 1

    def build(
        self,
        original_stream: BytesIO,
        extracted_texts: List[ExtractedText],
        redacted_texts: List[str],
    ) -> BytesIO:
        original_stream.seek(0)
        prs = Presentation(original_stream)

        idx_ref = [0]

        # Slides
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                self._replace_in_shape(shape, slide_idx, shape_idx, extracted_texts, redacted_texts, idx_ref)

            # Notas del slide
            if slide.has_notes_slide and slide.notes_slide:
                notes_text_frame = slide.notes_slide.notes_text_frame
                for para_idx, paragraph in enumerate(notes_text_frame.paragraphs):
                    for run_idx, run in enumerate(paragraph.runs):
                        idx = idx_ref[0]
                        if idx < len(extracted_texts):
                            loc = extracted_texts[idx].location
                            if (
                                loc[0] == "notes"
                                and loc[1] == slide_idx
                                and loc[2] == para_idx
                                and loc[3] == run_idx
                            ):
                                run.text = redacted_texts[idx]
                                idx_ref[0] = idx + 1

        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output
